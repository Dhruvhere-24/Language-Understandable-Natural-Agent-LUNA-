
def ocr_from_ppt(path: str) -> str:
    text = ""

    # ---- Step 1: Extract text from slides ----
    from pptx import Presentation
    prs = Presentation(path)

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text += " " + para.text

    # ---- Step 2: Extract images from PPT ----
    import zipfile
    from PIL import Image
    import io
    import pytesseract

    with zipfile.ZipFile(path, 'r') as zip_ref:
        for file in zip_ref.namelist():
            if file.startswith("ppt/media/"):
                image_data = zip_ref.read(file)

                image = Image.open(io.BytesIO(image_data))

                # OCR on image
                text += " " + pytesseract.image_to_string(image)

    return text

if __name__ == "__main__":
    import main

    file_path = "test2.pptx"
    sub_c, score = main.predict_subject(file_path)
    model = main.load_code_name_mapping("subjects_Final.json")
    sub_n = model.get(sub_c,"Unknown subject")

    print(f"\nsubject code is : {sub_c}\nsubject name is : {sub_n}\nconfidence score : {score}")

