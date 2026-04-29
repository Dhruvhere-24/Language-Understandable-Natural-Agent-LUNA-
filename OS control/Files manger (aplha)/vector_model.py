import os
import numpy as np
import pdfplumber
from docx import Document
from pptx import Presentation
from sentence_transformers import SentenceTransformer
from pdf_to_txt import smart_pdf_reader as spr, poppler_path # OCR for pdf
from ppt_to_text import ocr_from_ppt as ofp # OCR for ppt
from word_to_text import ocr_from_docx as ofd # OCR for word 


# Load embedding model
model = SentenceTransformer("all-MiniLM-L6-v2")

# Load subject embeddings
data = np.load("subject_embeddings.npz", allow_pickle=True)
subject_embeddings = data["embeddings"]
subject_codes = data["codes"]


# ---------- FILE TEXT EXTRACTION ----------

### Small problem here embedding model run 2 times
 
# def extract_text_from_pdf(path):
#     '''Simple txt extractor '''
#     text = ""
#     with pdfplumber.open(path) as pdf:
#         for page in pdf.pages:
#             page_text = page.extract_text()
#             if page_text:
#                 text += " " + page_text
#     return text


# def extract_text_from_docx(path):
#     doc = Document(path)
#     return " ".join([para.text for para in doc.paragraphs])


# def extract_text_from_pptx(path):
#     prs = Presentation(path)
#     text = ""
#     for slide in prs.slides:
#         for shape in slide.shapes:
#             if hasattr(shape, "text"):
#                 text += " " + shape.text
#     return text


def extract_text(file_path):
    ext = os.path.splitext(file_path)[1].lower()

    if ext == ".pdf":
        # return extract_text_from_pdf(file_path)
        return spr(file_path, poppler_path)
    elif ext == ".docx":
        return ofd(file_path)
    elif ext == ".pptx":
        return ofp(file_path)
    else:
        raise ValueError("Unsupported file type")

def load_code_name_mapping(json_path):
    import json
    with open(json_path, "r", encoding="utf-8") as f:
        subjects = json.load(f)

    return {s["subject_code"]: s["subject_name"] for s in subjects}


# ---------- CLASSIFICATION ----------

def predict_subject(file_path):
    print("Extracting text...")
    text = extract_text(file_path)

    if len(text.strip()) < 20:
        return "Extraction Failed", 0.0

    print("Generating embedding...")
    file_embedding = model.encode(text, normalize_embeddings=True)

    print("Computing similarity...")
    similarities = subject_embeddings @ file_embedding

    best_index = np.argmax(similarities)
    best_subject = subject_codes[best_index]
    confidence = similarities[best_index]

    return best_subject, float(confidence)


# ---------- TEST ----------

if __name__ == "__main__":
    file_path = "test.pdf" # path of file
    subject_code, score = predict_subject(file_path)

    # Load mapping once
    code_to_name = load_code_name_mapping("subjects_Final.json")
    subject_name = code_to_name.get(subject_code, "Unknown Subject")

    print("\nPredicted Subject Code:", subject_code)
    print("Predicted Subject Name:", subject_name)
    print("Confidence Score:", score)
